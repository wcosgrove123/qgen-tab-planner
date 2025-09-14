
from __future__ import annotations
from typing import Dict, Any, List
from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore

def build_deck(slides: List[Dict[str, Any]], outfile: str, theme):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for s in slides:
        slide = prs.slides.add_slide(blank)
        title = s.get("title") or ""
        subtitle = s.get("subtitle") or ""
        notes = s.get("notes") or ""
        chart_img = s.get("chart", {}).get("image_path")

        # Title
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11.5), Inches(0.6))
        tf = tbox.text_frame
        tf.text = title
        tf.paragraphs[0].font.name = theme.fonts["title"]["name"]
        tf.paragraphs[0].font.size = Pt(theme.fonts["title"]["size"])

        # Subtitle
        if subtitle:
            sbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(11.0), Inches(0.4))
            sf = sbox.text_frame
            sf.text = subtitle
            sf.paragraphs[0].font.name = theme.fonts["subtitle"]["name"]
            sf.paragraphs[0].font.size = Pt(theme.fonts["subtitle"]["size"])

        # Chart image
        if chart_img:
            slide.shapes.add_picture(chart_img, Inches(1), Inches(1.4), width=Inches(10))

        # Notes
        if notes:
            note = slide.notes_slide.notes_text_frame
            note.text = notes

    prs.save(outfile)
    return outfile
